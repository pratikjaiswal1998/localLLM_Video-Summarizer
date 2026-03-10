import os
import sys
import threading
import tempfile
import subprocess
import json
import re
import tkinter as tk
import customtkinter as ctk
from tkinter import messagebox
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

import yt_dlp
import ollama
from ollama import ResponseError

class SummarizerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Video Intelligence Summarizer")
        self.root.geometry("850x650")
        
        # Configure CustomTkinter Theme
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")
        self.root.grid_columnconfigure(0, weight=1)
        self.root.grid_rowconfigure(3, weight=1)
        
        # Auto-start Ollama in the background
        try:
            creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) if sys.platform == "win32" else 0
            app_path = os.path.expandvars(r"%LOCALAPPDATA%\Programs\Ollama\ollama app.exe")
            if os.path.exists(app_path):
                subprocess.Popen([app_path], creationflags=creation_flags, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            else:
                subprocess.Popen("ollama serve", shell=True, creationflags=creation_flags, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except Exception:
            pass
        
        # --- Top Input Area ---
        self.input_frame = ctk.CTkFrame(self.root, corner_radius=10)
        self.input_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.input_frame.grid_columnconfigure(1, weight=1)
        
        self.url_label = ctk.CTkLabel(self.input_frame, text="Video URL / File:", font=ctk.CTkFont(size=14, weight="bold"))
        self.url_label.grid(row=0, column=0, padx=15, pady=15)
        
        self.url_entry = ctk.CTkEntry(self.input_frame, placeholder_text="e.g. https://youtube.com/watch?v=...", font=ctk.CTkFont(size=13), height=35)
        self.url_entry.grid(row=0, column=1, padx=(0, 15), pady=15, sticky="ew")
        
        self.start_btn = ctk.CTkButton(self.input_frame, text="Summarize", font=ctk.CTkFont(size=14, weight="bold"), height=35, command=self.start_process)
        self.start_btn.grid(row=0, column=2, padx=(0, 15), pady=15)

        self.ppt_var = tk.BooleanVar(value=False)
        self.ppt_switch = ctk.CTkSwitch(self.input_frame, text="Generate PowerPoint (with Screenshots)", variable=self.ppt_var, font=ctk.CTkFont(size=12))
        self.ppt_switch.grid(row=1, column=0, columnspan=3, padx=15, pady=(0, 10), sticky="w")
        
        # --- Stage Tracker & Progress ---
        self.tracker_frame = ctk.CTkFrame(self.root, corner_radius=10, fg_color="transparent")
        self.tracker_frame.grid(row=1, column=0, padx=20, pady=(0, 10), sticky="ew")
        self.tracker_frame.grid_columnconfigure((0, 1, 2), weight=1)
        
        # 3-Stage Process Display
        self.stage_1_lbl = ctk.CTkLabel(self.tracker_frame, text="1. Extract Data", text_color="gray")
        self.stage_1_lbl.grid(row=0, column=0)
        self.stage_2_lbl = ctk.CTkLabel(self.tracker_frame, text="2. AI Listening", text_color="gray")
        self.stage_2_lbl.grid(row=0, column=1)
        self.stage_3_lbl = ctk.CTkLabel(self.tracker_frame, text="3. AI Writing", text_color="gray")
        self.stage_3_lbl.grid(row=0, column=2)
        
        # Progress Bar Layout
        self.progress = ctk.CTkProgressBar(self.root, height=8, corner_radius=4)
        self.progress.grid(row=2, column=0, padx=20, pady=(0, 10), sticky="ew")
        self.progress.set(0)
        
        # --- Console Output Terminal ---
        self.log_area = ctk.CTkTextbox(self.root, font=ctk.CTkFont(family="Consolas", size=12), corner_radius=10, border_width=1)
        self.log_area.grid(row=3, column=0, padx=20, pady=(0, 20), sticky="nsew")
        
    def log(self, message):
        self.log_area.insert(tk.END, message + "\n")
        self.log_area.see(tk.END)
        self.root.update_idletasks()
        
    def set_stage(self, stage_num):
        # Reset all
        self.stage_1_lbl.configure(text_color="gray", font=ctk.CTkFont(weight="normal"))
        self.stage_2_lbl.configure(text_color="gray", font=ctk.CTkFont(weight="normal"))
        self.stage_3_lbl.configure(text_color="gray", font=ctk.CTkFont(weight="normal"))
        
        if stage_num == 1:
            self.stage_1_lbl.configure(text_color="#4EC9B0", font=ctk.CTkFont(weight="bold"))
            self.progress.set(0.2)
        elif stage_num == 2:
            self.stage_1_lbl.configure(text_color="#4EC9B0")
            self.stage_2_lbl.configure(text_color="#4EC9B0", font=ctk.CTkFont(weight="bold"))
            self.progress.set(0.5)
        elif stage_num == 3:
            self.stage_1_lbl.configure(text_color="#4EC9B0")
            self.stage_2_lbl.configure(text_color="#4EC9B0")
            self.stage_3_lbl.configure(text_color="#4EC9B0", font=ctk.CTkFont(weight="bold"))
            self.progress.set(0.8)
        elif stage_num == 4: # Complete
            self.stage_1_lbl.configure(text_color="#4EC9B0")
            self.stage_2_lbl.configure(text_color="#4EC9B0")
            self.stage_3_lbl.configure(text_color="#4EC9B0")
            self.progress.set(1.0)

    def start_process(self):
        url = self.url_entry.get().strip()
        if not url:
            messagebox.showwarning("Warning", "Please enter a URL or file path.")
            return
            
        self.start_btn.configure(state="disabled")
        self.url_entry.configure(state="disabled")
        self.ppt_switch.configure(state="disabled")
        self.log_area.delete("0.0", tk.END)
        self.progress.set(0)
        self.set_stage(1)
        
        generate_ppt = self.ppt_var.get()
        
        # Run in background thread so GUI doesn't freeze
        threading.Thread(target=self.run_pipeline, args=(url, generate_ppt), daemon=True).start()
        
    def run_pipeline(self, url, generate_ppt):
        current_dir = os.path.dirname(os.path.abspath(sys.argv[0]))
        
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # --- Step 1: Download ---
                if url.startswith("http"):
                    self.log(f"[*] Downloading audio from {url}...")
                    
                    if generate_ppt:
                        self.log(f"[*] Mode: Video + Audio (for PowerPoint screenshots)...")
                        ydl_opts = {
                            'format': 'bestvideo[height<=720]+bestaudio/best',
                            'outtmpl': os.path.join(temp_dir, 'video.%(ext)s'),
                            'merge_output_format': 'mp4',
                            'quiet': True,
                            'no_warnings': True
                        }
                    else:
                        self.log(f"[*] Mode: Audio only (summarization)...")
                        ydl_opts = {
                            'format': 'bestaudio/best',
                            'outtmpl': os.path.join(temp_dir, '%(title)s.%(ext)s'),
                            'postprocessors': [{
                                'key': 'FFmpegExtractAudio',
                                'preferredcodec': 'mp3',
                                'preferredquality': '128',
                            }],
                            'quiet': True,
                            'no_warnings': True
                        }
                    
                    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                        info = ydl.extract_info(url, download=True)
                        filename = ydl.prepare_filename(info)
                        if generate_ppt:
                            audio_file = filename # It's a video file now, but whisper can read it
                        else:
                            name, _ = os.path.splitext(filename)
                            audio_file = name + ".mp3"
                else:
                    self.log(f"[*] Analyzing local file...")
                    audio_file = url
                    
                if not os.path.exists(audio_file):
                    self.log(f"[!] Error: Audio file {audio_file} not found.")
                    self.finish(success=False)
                    return
                    
                # --- Step 2: Transcribe ---
                self.set_stage(2)
                self.log(f"\n[*] Booting Local AI Listener (Faster-Whisper on 8GB VRAM)...")
                
                transcript_path = os.path.join(current_dir, "transcript.txt")
                transcriber_script = os.path.join(current_dir, "transcriber_backend.py")
                
                # Run transcriber in a subprocess so Windows forcefully reclaims the 3GB VRAM pool when it terminates
                cmd = [sys.executable, transcriber_script, audio_file, transcript_path]
                try:
                    creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) if sys.platform == "win32" else 0
                    process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, encoding="utf-8", creationflags=creation_flags)
                    
                    if process.stdout:
                        for line in iter(process.stdout.readline, ''):
                            if not line:
                                break
                            if line.strip():
                                self.log(line.strip('\n'))
                        
                    process.wait()
                    
                    if process.returncode != 0 or not os.path.exists(transcript_path):
                        self.log("[!] Transcription failed. See log above.")
                        self.finish(success=False)
                        return
                        
                    with open(transcript_path, "r", encoding="utf-8") as f:
                        transcript = f.read()

                    # Also load the timestamped version for PPT logic
                    timestamped_path = transcript_path.replace(".txt", "_timestamped.txt")
                    if os.path.exists(timestamped_path):
                        with open(timestamped_path, "r", encoding="utf-8") as f:
                            transcript_timestamped = f.read()
                    else:
                        transcript_timestamped = transcript
                        
                except Exception as e:
                    self.log(f"[!] Error running transcriber subprocess: {e}")
                    self.finish(success=False)
                    return
                    
                # --- VRAM CLEANUP ---
                self.log("\n[*] Whisper Subprocess Exited: Windows has forcefully reclaimed 100% of the VRAM for Llama 3.1!")
                
                # --- Step 3: Summarize ---
                self.set_stage(3)
                self.log(f"\n[*] Booting Local AI Writer (Llama 3.1 8B)...")
                
                # Make sure Ollama server is running before connecting
                import urllib.request
                import time
                
                def is_ollama_running():
                    try:
                        resp = urllib.request.urlopen("http://localhost:11434/", timeout=2)
                        return resp.getcode() == 200
                    except Exception:
                        return False

                if not is_ollama_running():
                    self.log("[*] Ollama server is not running. Starting it now in the background...")
                    try:
                        creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) if sys.platform == "win32" else 0
                        
                        # The CLI 'ollama serve' often fails silently on Windows. The most reliable way is to launch the actual System Tray app.
                        app_path = os.path.expandvars(r"%LOCALAPPDATA%\Programs\Ollama\ollama app.exe")
                        if os.path.exists(app_path):
                            subprocess.Popen([app_path], creationflags=creation_flags, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                        else:
                            subprocess.Popen("ollama serve", shell=True, creationflags=creation_flags, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                            
                    except Exception as boot_err:
                        self.log(f"[!] Could not start Ollama automatically: {boot_err}")
                        self.log("[HINT] Is Ollama installed? Please download it from https://ollama.com")
                    
                    # Wait up to 15 seconds for it to boot
                    started = False
                    for i in range(15):
                        if is_ollama_running():
                            started = True
                            break
                        time.sleep(1)
                    
                    if not started:
                        self.log("[!] Warning: Sent start command, but Ollama server didn't respond in time. Continuing anyway...")
                else:
                    self.log("[*] Ollama server is already running.")

                self.log(f"\n[*] Sending transcript to Ollama (Llama 3.1 8B)...")
                
                # Split transcript into chunks to prevent context window truncation on 1-hour videos
                def chunk_text(text, max_words=3000, overlap=100):
                    lines = text.split('\n')
                    chunks = []
                    current_chunk = []
                    current_words = 0
                    
                    for line in lines:
                        words_in_line = len(line.split())
                        if current_words + words_in_line > max_words:
                            chunks.append("\n".join(current_chunk))
                            # Keep some overlap lines
                            current_chunk = current_chunk[-5:] if len(current_chunk) > 5 else []
                            current_words = sum(len(l.split()) for l in current_chunk)
                        
                        current_chunk.append(line)
                        current_words += words_in_line
                    
                    if current_chunk:
                        chunks.append("\n".join(current_chunk))
                    return chunks

                # Use timestamped transcript for PPT chunking if enabled
                if generate_ppt:
                    transcript_chunks = chunk_text(transcript_timestamped)
                else:
                    transcript_chunks = chunk_text(transcript)
                    
                total_chunks = len(transcript_chunks)
                final_summary = ""
                all_slides_data = [] # Aggregate slides from all parts
                
                # Process each segment individually
                for idx, t_chunk in enumerate(transcript_chunks):
                    if total_chunks > 1:
                        self.log(f"[*] Writing Part {idx+1} of {total_chunks}...")
                        
                    if generate_ppt:
                        # Estimate how many slides this chunk should produce
                        chunk_words = len(t_chunk.split())
                        min_slides = max(3, chunk_words // 600)  # ~1 slide per 600 words, minimum 3

                        prompt = f"""Convert this transcript into exactly {min_slides} to {min_slides + 2} slides. Return ONLY a JSON array.

EXAMPLE of correct output with 4 slides:
[
  {{"title": "The Cost Of Bad Marketing", "bullets": ["Nike lost $25 billion in a single day after over-indexing on data-driven marketing", "Celebrity endorsements inflate costs, making companies structurally unprofitable", "Marketing managers optimize for measurable metrics to protect reputation, not growth", "Zerodha acquired millions of users at zero cost via YouTube trust-building"], "timestamp": 7.24}},
  {{"title": "Why Most Startups Fail", "bullets": ["80% of Indian startups are unprofitable due to chasing vanity metrics over fundamentals", "Raising hundreds of millions doesn't guarantee success — many funded startups go to zero", "The true secret is finding an untapped market, hiring talent, and executing consistently", "Survival is the strategy — companies that stick around long enough eventually win"], "timestamp": 120.5}},
  {{"title": "Understanding The Buyer", "bullets": ["Consumers watch 5-10 YouTube reviews before purchasing a phone above 20K", "The buying decision is emotional and gut-driven, not purely based on specs", "Brand recall drops sharply after 30 days without reinforcement", "Word of mouth and organic content drive more conversions than paid ads"], "timestamp": 250.0}},
  {{"title": "Building Without Capital", "bullets": ["AOS started with zero funding and grew through organic demand", "Money helps but is not the deciding factor in building large businesses", "Business case study analysts sensationalize failures for engagement", "Long-term consistency beats short-term capital injection every time"], "timestamp": 380.0}}
]

YOU MUST PRODUCE AT LEAST {min_slides} SLIDES. Each slide covers a DIFFERENT topic from the transcript.
- "title": 3-7 words, Title Case
- "bullets": array of 4-6 dense strings (15-25 words each, with names/numbers/facts)
- "timestamp": the [XX.XXs] second marker from the transcript where that topic starts

TRANSCRIPT:
{t_chunk}"""
                    else:
                        prompt = f"""
                        You are an expert analyst. Provide a detailed summary of the following transcript segment. 
                        Focus on logic, reasons why, and specific arguments.
                        Output in English.
                        
                        SEGMENT:
                        {t_chunk}
                        """
                    
                    def _call_ollama_ppt(messages):
                        """Call Ollama WITHOUT format='json' so the model can output naturally."""
                        try:
                            return ollama.chat(
                                model='llama3.1', messages=messages,
                                options={'num_ctx': 8192}
                            )
                        except Exception as e:
                            if "404" in str(e) or "not found" in str(e).lower():
                                self.log("[!] Llama 3.1 model not found. Downloading...")
                                for pd in ollama.pull('llama3.1', stream=True):
                                    pct = (pd.get('completed', 0) / pd.get('total', 1))
                                    self.progress.set(pct)
                                    self.root.update_idletasks()
                                self.log("[*] Download complete! Retrying...")
                                return ollama.chat(
                                    model='llama3.1', messages=messages,
                                    options={'num_ctx': 8192}
                                )
                            raise

                    def _extract_json_from_text(raw):
                        """Extract JSON array from free-form LLM output."""
                        s = raw.strip()
                        # Try markdown code block first
                        if "```" in s:
                            match = re.search(r"```(?:json)?\s*(.*?)\s*```", s, re.DOTALL)
                            if match: s = match.group(1).strip()
                        # Find the outermost array
                        start, end = s.find("["), s.rfind("]")
                        if start != -1 and end != -1 and end > start:
                            try:
                                return json.loads(s[start:end+1])
                            except json.JSONDecodeError:
                                pass
                        # Try as a single object or object with known wrapper keys
                        start, end = s.find("{"), s.rfind("}")
                        if start != -1 and end != -1 and end > start:
                            try:
                                obj = json.loads(s[start:end+1])
                                if isinstance(obj, dict):
                                    for key in ["slides", "data", "presentation", "slide"]:
                                        if key in obj and isinstance(obj[key], list):
                                            return obj[key]
                                    return [obj]
                            except json.JSONDecodeError:
                                # Try extracting multiple objects
                                objects = re.findall(r'\{[^{}]+\}', s)
                                if objects:
                                    try:
                                        return [json.loads(o) for o in objects]
                                    except json.JSONDecodeError:
                                        pass
                        return None

                    def _python_structure_transcript(chunk):
                        """Pure Python fallback: structure transcript into slides without LLM."""
                        # Parse timestamped lines
                        lines = []
                        for m in re.finditer(r'\[(\d+\.?\d*)s\s*->\s*\d+\.?\d*s\]\s*(.+)', chunk):
                            ts, text = float(m.group(1)), m.group(2).strip()
                            if len(text) > 15:  # Skip very short fragments
                                lines.append((ts, text))
                        if not lines:
                            # Non-timestamped: split by sentences
                            sentences = re.split(r'[.!?]+', chunk)
                            lines = [(i * 10, s.strip()) for i, s in enumerate(sentences) if len(s.strip()) > 15]

                        if not lines:
                            return []

                        # Group into slides of 5 lines each
                        slides = []
                        for i in range(0, len(lines), 5):
                            batch = lines[i:i+5]
                            bullets = [text for _, text in batch]
                            timestamp = batch[0][0]
                            # Auto-title from first bullet
                            words = bullets[0].replace(",", "").replace(".", "").split()
                            skip = {"the", "a", "an", "in", "on", "at", "for", "to", "and", "but", "or", "so",
                                    "is", "was", "are", "were", "it", "its", "this", "that", "you", "i", "we",
                                    "they", "he", "she", "if", "of", "with", "how", "what", "when", "where",
                                    "do", "don't", "just", "very", "really", "right", "like", "know", "think",
                                    "many", "most", "some", "lot", "lots", "people", "going", "have", "has"}
                            core = [w for w in words if w.lower() not in skip][:5]
                            if not core: core = words[:5]
                            title = " ".join(core).title()
                            slides.append({"title": title, "bullets": bullets, "timestamp": timestamp})

                        return slides

                    try:
                        if generate_ppt:
                            response = _call_ollama_ppt([
                                {'role': 'system', 'content': 'You are a slide generator. Output a JSON array and nothing else.'},
                                {'role': 'user', 'content': prompt}
                            ])
                        else:
                            response = ollama.chat(model='llama3.1', messages=[
                                {'role': 'user', 'content': prompt}
                            ], options={'num_ctx': 8192})
                    except Exception as e:
                        self.log(f"[!] Ollama Error: {e}")
                        continue

                    if total_chunks > 1:
                        final_summary += f"\n\n--- PART {idx+1} SUMMARY ---\n"

                    content = response['message']['content']

                    if generate_ppt:
                        try:
                            slides_from_llm = []
                            parsed = _extract_json_from_text(content)

                            if parsed and isinstance(parsed, list):
                                # Check if we got proper slide objects
                                for item in parsed:
                                    if isinstance(item, dict):
                                        bullets = item.get("bullets", item.get("bullet_points", item.get("points", [])))
                                        if isinstance(bullets, str): bullets = [bullets]
                                        if isinstance(bullets, list) and bullets:
                                            slides_from_llm.append({
                                                "title": item.get("title", item.get("heading", item.get("header", "Key Points"))),
                                                "bullets": bullets,
                                                "timestamp": self.clean_timestamp(item.get("timestamp", item.get("time", 0)))
                                            })
                                        else:
                                            # Dict without bullets — salvage string values
                                            vals = [str(v) for v in item.values() if isinstance(v, str) and len(str(v)) > 15]
                                            if vals:
                                                slides_from_llm.append({
                                                    "title": str(list(item.keys())[0]).title() if item else "Key Points",
                                                    "bullets": vals,
                                                    "timestamp": 0
                                                })

                            if slides_from_llm:
                                self.log(f"[*] LLM produced {len(slides_from_llm)} structured slides")
                                for sl in slides_from_llm:
                                    all_slides_data.append(sl)
                                    final_summary += f"### {sl['title']}\n" + "\n".join(f"- {b}" for b in sl['bullets']) + "\n\n"

                                # If LLM produced too few slides, supplement with Python structuring
                                if len(slides_from_llm) < min_slides:
                                    # Find the last timestamp covered by LLM slides
                                    last_covered_ts = max(sl.get("timestamp", 0) for sl in slides_from_llm)
                                    # Extract remaining transcript lines after that timestamp
                                    remaining_lines = []
                                    for m in re.finditer(r'\[(\d+\.?\d*)s\s*->\s*\d+\.?\d*s\]\s*(.+)', t_chunk):
                                        ts, text = float(m.group(1)), m.group(2).strip()
                                        if ts > last_covered_ts + 30 and len(text) > 15:
                                            remaining_lines.append((ts, text))

                                    if remaining_lines:
                                        needed = min_slides - len(slides_from_llm)
                                        lines_per_slide = max(4, len(remaining_lines) // max(needed, 1))
                                        self.log(f"[*] Supplementing with {needed} Python-structured slides from remaining transcript")
                                        for si in range(0, len(remaining_lines), lines_per_slide):
                                            batch = remaining_lines[si:si + lines_per_slide]
                                            if not batch: break
                                            bullets = [text for _, text in batch[:6]]
                                            ts = batch[0][0]
                                            words = bullets[0].replace(",", "").replace(".", "").split()
                                            skip_words = {"the","a","an","in","on","at","for","to","and","but","or","so","is","was","are","were","it","its","this","that","you","i","we","they","he","she","if","of","with","how","what","when","where","do","don't","just","very","really","right","like","know","think"}
                                            core = [w for w in words if w.lower() not in skip_words][:5]
                                            if not core: core = words[:5]
                                            title = " ".join(core).title()
                                            sl = {"title": title, "bullets": bullets, "timestamp": ts}
                                            all_slides_data.append(sl)
                                            final_summary += f"### {title}\n" + "\n".join(f"- {b}" for b in bullets) + "\n\n"
                            else:
                                # LLM failed completely — full Python structuring
                                self.log(f"[!] LLM output not usable — using Python transcript structurer")
                                python_slides = _python_structure_transcript(t_chunk)
                                if python_slides:
                                    self.log(f"[*] Python structurer produced {len(python_slides)} slides from transcript")
                                    for sl in python_slides:
                                        all_slides_data.append(sl)
                                        final_summary += f"### {sl['title']}\n" + "\n".join(f"- {b}" for b in sl['bullets']) + "\n\n"
                                else:
                                    self.log(f"[!] Could not structure chunk {idx+1} — skipping")

                        except Exception as parse_err:
                            self.log(f"[!] Chunk {idx+1} parse error ({type(parse_err).__name__}: {parse_err})")
                            # Last resort: Python structuring
                            python_slides = _python_structure_transcript(t_chunk)
                            for sl in python_slides:
                                all_slides_data.append(sl)
                                final_summary += f"### {sl['title']}\n" + "\n".join(f"- {b}" for b in sl['bullets']) + "\n\n"
                    else:
                        final_summary += content
                
                # FINAL STEP: Build the PPT once from all aggregated data
                if generate_ppt and all_slides_data:
                    self.log(f"[*] Building Final PowerPoint with {len(all_slides_data)} slides...")
                    try:
                        self.create_ppt(all_slides_data, audio_file, current_dir)
                        final_summary += "[*] PowerPoint presentation generated successfully!\n"
                    except Exception as e:
                        self.log(f"[!] Critical Error building final PPT: {e}")

                # Purge Llama 3.1 from VRAM immediately so next Whisper run has a clean slate
                try:
                    ollama.chat(model='llama3.1', messages=[{'role': 'user', 'content': ''}], keep_alive=0)
                    self.log("[*] VRAM purged: Llama 3.1 unloaded from GPU.")
                except Exception:
                    pass
                
                summary = final_summary.strip()
                
                summary_path = os.path.join(current_dir, "summary.txt")
                with open(summary_path, "w", encoding="utf-8") as f:
                    f.write(summary)
                
                self.log("\n" + "="*50)
                self.log("SUMMARY OF VIDEO:")
                self.log("="*50)
                self.log(summary)
                self.log("="*50)
                self.log(f"\n[*] Saved summary to {summary_path}")
                
                # Auto-open the results
                try:
                    if generate_ppt:
                        ppt_file = os.path.join(current_dir, "presentation.pptx")
                        if os.path.exists(ppt_file):
                            os.startfile(ppt_file)
                    else:
                        os.startfile(summary_path)
                except Exception:
                    pass
                
                self.finish(success=True)
                
        except Exception as e:
            self.log(f"\n[!] An error occurred: {str(e)}")
            if "ollama" in str(e).lower():
                self.log("\n[HINT] Make sure Ollama is installed and you ran: 'ollama run llama3.1' in command prompt!")
            self.finish(success=False)

    def clean_timestamp(self, ts):
        """Converts various timestamp formats (seconds, MM:SS, HH:MM:SS) to float seconds."""
        if isinstance(ts, (int, float)): return float(ts)
        if isinstance(ts, str):
            ts = ts.strip().replace("[", "").replace("]", "")
            if ":" in ts:
                parts = ts.split(':')
                try:
                    if len(parts) == 2: return int(parts[0]) * 60 + float(parts[1])
                    if len(parts) == 3: return int(parts[0]) * 3600 + int(parts[1]) * 60 + float(parts[2])
                except (ValueError, TypeError): return 0.0
            try: return float(ts)
            except (ValueError, TypeError): return 0.0
        return 0.0

    # --- Entity name corrections for common transcription errors ---
    ENTITY_CORRECTIONS = {
        "Anthropik": "Anthropic",
        "anthropik": "Anthropic",
        "ANTHROPIK": "ANTHROPIC",
        "Amadei": "Amodei",
        "amadei": "amodei",
        "AMADEI": "AMODEI",
        "Dario Amadei": "Dario Amodei",
    }

    def _fix_entities(self, text):
        """Fix common transcription misspellings of known entity names."""
        for wrong, right in self.ENTITY_CORRECTIONS.items():
            text = text.replace(wrong, right)
        return text

    def _add_slide_background(self, slide, prs, shapes_module, bg_color, accent_color):
        """Add dark background + accent bar to a slide."""
        bg_rect = slide.shapes.add_shape(shapes_module.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = bg_color
        bg_rect.line.fill.background()
        # Accent bar — thicker (0.25") so it's visible when projected
        accent_bar = slide.shapes.add_shape(shapes_module.RECTANGLE, 0, 0, prs.slide_width, Inches(0.25))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()

    def _add_slide_number(self, slide, prs, slide_num, total, muted_color):
        """Add a slide number in the bottom-right corner."""
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
        tf = num_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {total}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = 'Segoe UI'
        p.font.size = Pt(11)
        p.font.color.rgb = muted_color

    def create_ppt(self, slides_data, video_file, output_dir):
        """Builds a polished, dark-themed PowerPoint deck with title slide, varied layouts, and proper typography."""
        self.log("[*] Designing PowerPoint Slides...")

        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn

        prs = Presentation()

        # 16:9 Widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Theme Colors
        BG_COLOR = RGBColor(30, 30, 30)        # Dark Charcoal
        ACCENT_COLOR = RGBColor(0, 122, 204)   # Blue
        TITLE_COLOR = RGBColor(255, 255, 255)   # Pure white for titles
        TEXT_COLOR = RGBColor(220, 220, 220)    # Off-white for body
        MUTED_COLOR = RGBColor(130, 130, 130)   # Muted gray for captions/numbers
        BULLET_COLOR = ACCENT_COLOR              # Blue bullet markers

        # Font with fallback
        FONT = 'Segoe UI'

        # Total slides = title + content + closing
        total_slide_count = len(slides_data) + 2

        with tempfile.TemporaryDirectory() as screenshot_dir:

            # ============================================================
            # TITLE SLIDE
            # ============================================================
            slide_layout = prs.slide_layouts[6]
            title_slide = prs.slides.add_slide(slide_layout)
            self._add_slide_background(title_slide, prs, MSO_SHAPE, BG_COLOR, ACCENT_COLOR)

            # Big centered title
            first_title = self._fix_entities(slides_data[0].get("title", "Video Summary"))
            title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = first_title.upper()
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR

            # Subtitle line
            sub_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = "AI-Generated Video Summary"
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.font.name = FONT
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = MUTED_COLOR

            # Accent divider line under subtitle
            divider = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(5.0), Inches(3.3), Inches(0.04))
            divider.fill.solid()
            divider.fill.fore_color.rgb = ACCENT_COLOR
            divider.line.fill.background()

            self._add_slide_number(title_slide, prs, 1, total_slide_count, MUTED_COLOR)

            # ============================================================
            # CONTENT SLIDES
            # ============================================================
            for i, slide_info in enumerate(slides_data):
                if not isinstance(slide_info, dict): continue

                title_text = self._fix_entities(slide_info.get("title", f"Slide {i+1}"))
                bullets = slide_info.get("bullets", [])

                if isinstance(bullets, str): bullets = [bullets]
                elif not isinstance(bullets, list): bullets = []

                # Fix entity names in bullets
                bullets = [self._fix_entities(b) for b in bullets]

                raw_ts = slide_info.get("timestamp", 0)
                timestamp = self.clean_timestamp(raw_ts)

                # 1. Capture Screenshot
                screenshot_path = os.path.join(screenshot_dir, f"slide_{i}.jpg")
                has_screenshot = False
                try:
                    h = int(timestamp // 3600)
                    m = int((timestamp % 3600) // 60)
                    s = timestamp % 60
                    ts_str = f"{h:02d}:{m:02d}:{s:05.2f}"

                    creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) if sys.platform == "win32" else 0
                    cmd = [
                        'ffmpeg', '-y', '-ss', ts_str, '-i', video_file,
                        '-vframes', '1', '-q:v', '2', screenshot_path
                    ]
                    subprocess.run(cmd, creationflags=creation_flags, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    has_screenshot = (os.path.exists(screenshot_path) and os.path.getsize(screenshot_path) > 1000)
                except Exception:
                    has_screenshot = False

                # 2. Add Slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_slide_background(slide, prs, MSO_SHAPE, BG_COLOR, ACCENT_COLOR)

                # 3. Title — dynamic font sizing with word wrap
                title_display = title_text.upper()
                title_font_size = Pt(40) if len(title_display) <= 45 else Pt(34) if len(title_display) <= 60 else Pt(28)

                title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.2))
                tf_title = title_box.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.paragraphs[0]
                p_title.text = title_display
                p_title.font.name = FONT
                p_title.font.size = title_font_size
                p_title.font.bold = True
                p_title.font.color.rgb = TITLE_COLOR

                if has_screenshot:
                    # --- TWO-COLUMN LAYOUT: image left, bullets right ---

                    # Clamp image height to available space
                    img_left = Inches(0.6)
                    img_top = Inches(1.9)
                    img_max_width = Inches(6.3)
                    img_max_height = Inches(4.8)

                    try:
                        from PIL import Image as PILImage
                        with PILImage.open(screenshot_path) as img:
                            iw, ih = img.size
                        aspect = iw / ih
                        calc_height = img_max_width / aspect
                        if calc_height > img_max_height:
                            pic = slide.shapes.add_picture(screenshot_path, img_left, img_top, height=img_max_height)
                        else:
                            pic = slide.shapes.add_picture(screenshot_path, img_left, img_top, width=img_max_width)
                    except ImportError:
                        pic = slide.shapes.add_picture(screenshot_path, img_left, img_top, width=img_max_width)

                    # Image border (1pt light gray outline)
                    pic.line.color.rgb = RGBColor(80, 80, 80)
                    pic.line.width = Pt(1)

                    # Timestamp caption below image
                    cap_top = img_top + pic.height + Inches(0.1)
                    cap_box = slide.shapes.add_textbox(img_left, cap_top, img_max_width, Inches(0.3))
                    tf_cap = cap_box.text_frame
                    p_cap = tf_cap.paragraphs[0]
                    mins = int(timestamp // 60)
                    secs = int(timestamp % 60)
                    p_cap.text = f"Timestamp: {mins}:{secs:02d}"
                    p_cap.font.name = FONT
                    p_cap.font.size = Pt(11)
                    p_cap.font.color.rgb = MUTED_COLOR
                    p_cap.font.italic = True

                    # Bullets on right side
                    bullet_left = Inches(7.3)
                    bullet_width = Inches(5.4)
                else:
                    # --- FULL-WIDTH LAYOUT: no screenshot ---
                    self.log(f"[!] Screenshot fail for slide {i+1} at {timestamp}s (using full-width layout)")
                    bullet_left = Inches(0.8)
                    bullet_width = Inches(11.7)

                # 4. Bullet Points — dynamic font + spacing based on count
                num_bullets = len(bullets)
                if has_screenshot:
                    if num_bullets <= 2: bullet_top = Inches(2.8)
                    elif num_bullets <= 4: bullet_top = Inches(2.0)
                    else: bullet_top = Inches(1.9)
                else:
                    if num_bullets <= 3: bullet_top = Inches(2.4)
                    else: bullet_top = Inches(2.0)

                # Font scales down as bullets increase; spacing tightens proportionally
                if num_bullets <= 3:
                    bullet_font_size, bullet_spacing = Pt(22), Pt(14)
                elif num_bullets <= 5:
                    bullet_font_size, bullet_spacing = Pt(18), Pt(10)
                elif num_bullets <= 7:
                    bullet_font_size, bullet_spacing = Pt(16), Pt(8)
                else:
                    bullet_font_size, bullet_spacing = Pt(14), Pt(6)

                bullet_height = Inches(5.2)

                txBox = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
                tf = txBox.text_frame
                tf.word_wrap = True

                for bi, b in enumerate(bullets):
                    if bi == 0:
                        p = tf.paragraphs[0]  # Use the existing first paragraph — no empty gap
                    else:
                        p = tf.add_paragraph()
                    p.space_before = bullet_spacing
                    p.level = 0

                    # Add colored bullet character as a separate run, then body text
                    bullet_run = p.add_run()
                    bullet_run.text = "\u2022  "
                    bullet_run.font.name = FONT
                    bullet_run.font.size = bullet_font_size
                    bullet_run.font.color.rgb = BULLET_COLOR

                    text_run = p.add_run()
                    text_run.text = b
                    text_run.font.name = FONT
                    text_run.font.size = bullet_font_size
                    text_run.font.color.rgb = TEXT_COLOR

                # 5. Slide number
                self._add_slide_number(slide, prs, i + 2, total_slide_count, MUTED_COLOR)

            # ============================================================
            # CLOSING SLIDE
            # ============================================================
            closing_slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_background(closing_slide, prs, MSO_SHAPE, BG_COLOR, ACCENT_COLOR)

            # "Key Takeaways" header
            close_title_box = closing_slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2))
            tf_ct = close_title_box.text_frame
            tf_ct.word_wrap = True
            p_ct = tf_ct.paragraphs[0]
            p_ct.text = "KEY TAKEAWAYS"
            p_ct.alignment = PP_ALIGN.CENTER
            p_ct.font.name = FONT
            p_ct.font.size = Pt(44)
            p_ct.font.bold = True
            p_ct.font.color.rgb = TITLE_COLOR

            # Pull first bullet from each content slide as a recap
            # Collect recap items first so we can size dynamically
            recap_items = []
            for sd in slides_data:
                if not isinstance(sd, dict): continue
                recap_bullets = sd.get("bullets", [])
                if isinstance(recap_bullets, str): recap_bullets = [recap_bullets]
                first_point = self._fix_entities(recap_bullets[0]) if recap_bullets else ""
                if first_point:
                    recap_items.append(first_point)

            # Dynamic font + spacing based on recap count
            recap_count = len(recap_items)
            if recap_count <= 5:
                recap_font, recap_spacing, recap_top, recap_height = Pt(20), Pt(12), Inches(2.8), Inches(3.8)
            elif recap_count <= 8:
                recap_font, recap_spacing, recap_top, recap_height = Pt(16), Pt(8), Inches(2.6), Inches(4.2)
            elif recap_count <= 12:
                recap_font, recap_spacing, recap_top, recap_height = Pt(13), Pt(6), Inches(2.4), Inches(4.4)
            else:
                recap_font, recap_spacing, recap_top, recap_height = Pt(11), Pt(4), Inches(2.2), Inches(4.6)

            recap_box = closing_slide.shapes.add_textbox(Inches(1.0), recap_top, Inches(11.3), recap_height)
            tf_recap = recap_box.text_frame
            tf_recap.word_wrap = True

            for ri, point in enumerate(recap_items):
                if ri == 0:
                    p = tf_recap.paragraphs[0]
                else:
                    p = tf_recap.add_paragraph()
                p.space_before = recap_spacing

                br = p.add_run()
                br.text = "\u2022  "
                br.font.name = FONT
                br.font.size = recap_font
                br.font.color.rgb = ACCENT_COLOR

                tr = p.add_run()
                tr.text = point
                tr.font.name = FONT
                tr.font.size = recap_font
                tr.font.color.rgb = TEXT_COLOR

            # Disclaimer
            disc_box = closing_slide.shapes.add_textbox(Inches(2.0), Inches(6.5), Inches(9.3), Inches(0.5))
            tf_disc = disc_box.text_frame
            p_disc = tf_disc.paragraphs[0]
            p_disc.text = "Auto-generated summary \u2014 verify claims against original source before citing."
            p_disc.alignment = PP_ALIGN.CENTER
            p_disc.font.name = FONT
            p_disc.font.size = Pt(11)
            p_disc.font.italic = True
            p_disc.font.color.rgb = MUTED_COLOR

            self._add_slide_number(closing_slide, prs, total_slide_count, total_slide_count, MUTED_COLOR)

        ppt_path = os.path.join(output_dir, "presentation.pptx")
        prs.save(ppt_path)
        self.log(f"[*] PowerPoint saved to: {ppt_path}")


    def finish(self, success):
        if success:
            self.set_stage(4)
        else:
            self.stage_1_lbl.configure(text_color="red")
            self.stage_2_lbl.configure(text_color="red")
            self.stage_3_lbl.configure(text_color="red")
            self.progress.set(1.0)
            self.progress.configure(progress_color="red")
            
        self.start_btn.configure(state="normal")
        self.url_entry.configure(state="normal")
        self.ppt_switch.configure(state="normal")

if __name__ == "__main__":
    root = ctk.CTk()
    app = SummarizerApp(root)
    root.mainloop()
